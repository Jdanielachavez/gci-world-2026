"""Build the Final Assignment business-proposal deck (<=15 slides) as PPTX, then export to PDF."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ASSET = 'slide_assets/'

# ---- Palette ----
NAVY   = RGBColor(0x0B, 0x25, 0x45)
BLUE   = RGBColor(0x2C, 0x7D, 0xA0)
TEAL   = RGBColor(0x46, 0x8F, 0xAF)
RED    = RGBColor(0xE6, 0x39, 0x46)
GREEN  = RGBColor(0x2A, 0x9D, 0x8F)
LIGHT  = RGBColor(0xE8, 0xF1, 0xF5)
GRAY   = RGBColor(0x6C, 0x75, 0x7D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x21, 0x25, 0x29)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font='Calibri', italic=False, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb

def bullets(s, x, y, w, h, items, size=16, color=DARK, gap=6, bullet='▪  ', bold_lead=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        if isinstance(it, tuple):
            lead, rest = it
            r0 = p.add_run(); r0.text = bullet + lead
            r0.font.size = Pt(size); r0.font.bold = True; r0.font.color.rgb = color; r0.font.name = 'Calibri'
            r1 = p.add_run(); r1.text = rest
            r1.font.size = Pt(size); r1.font.bold = False; r1.font.color.rgb = color; r1.font.name = 'Calibri'
        else:
            r = p.add_run(); r.text = bullet + it
            r.font.size = Pt(size); r.font.bold = bold_lead; r.font.color.rgb = color; r.font.name = 'Calibri'
    return tb

def header(s, n, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Pt(4), RED)
    if kicker:
        txt(s, Inches(0.55), Inches(0.13), Inches(11), Inches(0.3), kicker.upper(),
            size=12, color=LIGHT, bold=True)
    txt(s, Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.7), title,
        size=26, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # page number
    txt(s, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35), str(n),
        size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def footer_src(s, text):
    txt(s, Inches(0.55), Inches(7.02), Inches(11.5), Inches(0.4), text,
        size=9.5, color=GRAY, italic=True)

def pic(s, path, x, y, w=None, h=None):
    if w and h: return s.shapes.add_picture(path, x, y, w, h)
    if w: return s.shapes.add_picture(path, x, y, width=w)
    return s.shapes.add_picture(path, x, y, height=h)

def kpi_card(s, x, y, w, h, value, label, vcolor=NAVY):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, w, Pt(5), RED)
    txt(s, x, y+Inches(0.18), w, Inches(0.6), value, size=30, color=vcolor, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, y+h-Inches(0.62), w, Inches(0.55), label, size=12.5, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1 — Title
# ============================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Pt(4), RED)
rect(s, Inches(0.0), Inches(0.0), Inches(0.25), SH, RED)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5), "GCI WORLD 2026 APRIL  ·  FINAL ASSIGNMENT  ·  PoC PROPOSAL",
    size=15, color=TEAL, bold=True)
txt(s, Inches(0.9), Inches(2.4), Inches(11.7), Inches(1.8),
    "Turning Customer Data into Retained Revenue",
    size=44, color=WHITE, bold=True, line_spacing=1.0)
txt(s, Inches(0.9), Inches(3.95), Inches(11.5), Inches(0.6),
    "A machine-learning churn-prevention proposal for Company A (Telecom)",
    size=20, color=LIGHT)
txt(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.5),
    "Prepared by:  Daniela Chavez", size=16, color=WHITE, bold=True)
txt(s, Inches(0.9), Inches(5.45), Inches(11.5), Inches(0.5),
    "IT Consulting — Data Science Associate", size=14, color=GRAY)
txt(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
    "Matsuo–Iwasawa Laboratory, The University of Tokyo", size=12, color=GRAY, italic=True)

# ============================================================
# SLIDE 2 — Executive Summary
# ============================================================
s = slide(); header(s, 2, "Executive Summary", "The bottom line first")
txt(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.9),
    "Company A loses nearly half its customers. A data-driven, ML-targeted retention program "
    "can recover ~$6M in revenue at almost no incremental cost.",
    size=18, color=NAVY, bold=True, line_spacing=1.1)
kpi_card(s, Inches(0.55), Inches(2.7), Inches(2.85), Inches(1.7), "49.6%", "Customer churn rate", RED)
kpi_card(s, Inches(3.65), Inches(2.7), Inches(2.85), Inches(1.7), "0.69", "Model ROC-AUC", GREEN)
kpi_card(s, Inches(6.75), Inches(2.7), Inches(2.85), Inches(1.7), "1.67x", "Lift in top 5% at-risk", BLUE)
kpi_card(s, Inches(9.85), Inches(2.7), Inches(2.85), Inches(1.7), "$6.0M", "Net benefit (24-mo)", NAVY)
bullets(s, Inches(0.55), Inches(4.8), Inches(12.2), Inches(2.2), [
    ("Problem:  ", "≈49.6% of ~100,000 customers churn within the 31–60 day window — a major recurring revenue leak."),
    ("Solution:  ", "A monthly churn-risk scoring engine ranks customers so the retention team targets the highest-risk segment first."),
    ("Evidence:  ", "Gradient Boosting (ROC-AUC 0.69) identifies churners 1.5–1.7× better than random targeting."),
    ("Impact:  ", "Targeting the top 20% yields ~$6.0M net benefit and ~$2.0M more than untargeted outreach."),
], size=15, gap=7)

# ============================================================
# SLIDE 3 — Market Analysis
# ============================================================
s = slide(); header(s, 3, "Market Context: Why Churn Is the Battleground", "Industry analysis")
bullets(s, Inches(0.55), Inches(1.5), Inches(6.4), Inches(5), [
    ("US telecom is a ~$650B market ", "[1] — mature, saturated, and fiercely competitive."),
    ("Annual churn averages 20–30% ", "across wireless carriers [2]; acquisition is largely zero-sum."),
    ("Acquiring a new customer costs 5–7× ", "more than retaining an existing one [3]."),
    ("A 5% lift in retention can raise profits 25–95% ", "(Reichheld & Sasser) [4]."),
    ("Implication: ", "with little room to grow the pie, defending the existing base is the highest-ROI lever available to Company A."),
], size=15, gap=11)
# right callout panel
rect(s, Inches(7.35), Inches(1.6), Inches(5.4), Inches(4.6), LIGHT)
rect(s, Inches(7.35), Inches(1.6), Inches(5.4), Pt(5), RED)
txt(s, Inches(7.7), Inches(1.9), Inches(4.8), Inches(0.5), "THE RETENTION ECONOMICS", size=14, color=NAVY, bold=True)
txt(s, Inches(7.7), Inches(2.55), Inches(4.8), Inches(1.2),
    "5–7×", size=44, color=RED, bold=True)
txt(s, Inches(7.7), Inches(3.5), Inches(4.8), Inches(0.6),
    "cost to acquire vs. retain a customer", size=14, color=DARK)
txt(s, Inches(7.7), Inches(4.3), Inches(4.8), Inches(1.2),
    "25–95%", size=40, color=GREEN, bold=True)
txt(s, Inches(7.7), Inches(5.2), Inches(4.8), Inches(0.8),
    "potential profit increase from a\n5% improvement in retention", size=14, color=DARK)
footer_src(s, "Sources: [1] IBISWorld/Statista 2024  ·  [2] McKinsey Telecom  ·  [3] HBR  ·  [4] Bain & Co. — full list on References slide.")

# ============================================================
# SLIDE 4 — The Dataset
# ============================================================
s = slide(); header(s, 4, "The Data We Were Given", "Proof of Concept dataset")
bullets(s, Inches(0.55), Inches(1.55), Inches(6.3), Inches(4.8), [
    ("Two tables, joined on Customer_ID:", ""),
    ("  • Client.csv ", "— 100,000 customers × 50 demographic/account columns."),
    ("  • Record.csv ", "— 100,000 rows × 51 usage/behavioral columns."),
    ("Merged view: ", "100,000 customers × 100 features."),
    ("Target variable: ", "churn (1 = left within 31–60 days, 0 = stayed)."),
    ("Real-world messiness: ", "many columns have missing values and ambiguous definitions — handled explicitly in preprocessing."),
], size=15, gap=9)
# feature family panel
rect(s, Inches(7.25), Inches(1.55), Inches(5.5), Inches(4.9), LIGHT)
txt(s, Inches(7.55), Inches(1.75), Inches(5), Inches(0.4), "FEATURE FAMILIES", size=14, color=NAVY, bold=True)
bullets(s, Inches(7.55), Inches(2.35), Inches(5, ), Inches(4), [
    ("Usage: ", "minutes of use, calls placed/received, peak vs off-peak"),
    ("Revenue: ", "monthly charge, overage, billing-adjusted totals"),
    ("Service quality: ", "dropped / blocked / unanswered calls"),
    ("Equipment: ", "handset price, equipment age, web capability"),
    ("Account: ", "tenure (months), # subscribers, credit class"),
    ("Demographics: ", "income, region, dwelling, household"),
], size=13.5, gap=8)
footer_src(s, "Source: Company A PoC dataset (anonymized telecom provider), provided in course materials.")

# ============================================================
# SLIDE 5 — EDA: churn scale
# ============================================================
s = slide(); header(s, 5, "Finding 1: Half the Base Is Walking Out", "Exploratory data analysis")
pic(s, ASSET+'01_churn_dist.png', Inches(0.5), Inches(1.55), w=Inches(8.0))
rect(s, Inches(8.85), Inches(1.7), Inches(3.95), Inches(4.6), LIGHT)
rect(s, Inches(8.85), Inches(1.7), Inches(3.95), Pt(5), RED)
txt(s, Inches(9.1), Inches(2.0), Inches(3.5), Inches(0.5), "SO WHAT?", size=14, color=NAVY, bold=True)
bullets(s, Inches(9.1), Inches(2.6), Inches(3.45), Inches(3.6), [
    "49.6% churn is extreme — almost a coin-flip on every customer.",
    "Even a small, well-targeted reduction protects millions in recurring revenue.",
    "The base is balanced, so the modeling challenge is ranking risk, not finding a rare event.",
], size=14, gap=12)
footer_src(s, "Analysis of merged Company A dataset (n = 100,000).")

# ============================================================
# SLIDE 6 — EDA: drivers
# ============================================================
s = slide(); header(s, 6, "Finding 2: What Separates Churners From Stayers", "Exploratory data analysis")
pic(s, ASSET+'02_correlations.png', Inches(0.45), Inches(1.5), w=Inches(6.6))
pic(s, ASSET+'03_drivers.png', Inches(7.2), Inches(1.5), w=Inches(5.9))
txt(s, Inches(7.2), Inches(5.0), Inches(5.9), Inches(1.7),
    "Equipment age and handset price are the strongest signals: customers on older, "
    "cheaper phones — and those with declining usage — are the most likely to leave. "
    "These map directly to actionable retention offers (device upgrades, loyalty perks).",
    size=13.5, color=DARK, line_spacing=1.1)
footer_src(s, "Left: top correlations with churn. Right: distributions of equipment age and tenure by churn status.")

# ============================================================
# SLIDE 7 — Problem definition
# ============================================================
s = slide(); header(s, 7, "Problem Definition & ML Task", "From insight to objective")
rect(s, Inches(0.55), Inches(1.55), Inches(5.9), Inches(2.2), LIGHT)
txt(s, Inches(0.8), Inches(1.75), Inches(5.4), Inches(0.5), "BUSINESS OBJECTIVE", size=14, color=NAVY, bold=True)
txt(s, Inches(0.8), Inches(2.35), Inches(5.4), Inches(1.3),
    "Reduce revenue lost to churn by identifying — in advance — which customers are about "
    "to leave, so retention can act before they do.", size=15.5, color=DARK, line_spacing=1.1)
rect(s, Inches(6.75), Inches(1.55), Inches(6.0), Inches(2.2), LIGHT)
txt(s, Inches(7.0), Inches(1.75), Inches(5.4), Inches(0.5), "MACHINE-LEARNING TASK", size=14, color=NAVY, bold=True)
txt(s, Inches(7.0), Inches(2.35), Inches(5.5), Inches(1.3),
    "Binary classification of churn, optimized to RANK customers by risk probability rather "
    "than to label them yes/no.", size=15.5, color=DARK, line_spacing=1.1)
txt(s, Inches(0.55), Inches(4.1), Inches(12), Inches(0.5), "Why this framing?", size=17, color=NAVY, bold=True)
bullets(s, Inches(0.55), Inches(4.65), Inches(12.2), Inches(2.2), [
    ("Target = churn ", "— directly ties to the client's core profitability problem (validated in EDA)."),
    ("Metric = ROC-AUC ", "— measures how well we RANK risk; insensitive to threshold, ideal when the plan is to contact a top-N list."),
    ("Ranking > hard labels ", "— the retention budget is finite, so we prioritize the customers most likely to churn AND worth saving."),
], size=15, gap=9)

# ============================================================
# SLIDE 8 — Modeling approach
# ============================================================
s = slide(); header(s, 8, "Modeling Approach", "How we built it")
bullets(s, Inches(0.55), Inches(1.5), Inches(5.7), Inches(5), [
    ("Preprocessing pipeline (leak-free):", ""),
    ("  • Numeric ", "— median imputation + standardization."),
    ("  • Categorical ", "— most-frequent imputation + one-hot encoding."),
    ("  • Dropped ", "high-cardinality / identifier columns."),
    ("Feature engineering:", ""),
    ("  • failure_ratio ", "— dropped+blocked ÷ attempted calls (service quality)."),
    ("  • rev_per_minute ", "— revenue efficiency."),
    ("  • custcare_intensity ", "— support-call rate (dissatisfaction proxy)."),
    ("Validation: ", "stratified 80/20 split + 5-fold cross-validation, fixed random seed (reproducible)."),
], size=14, gap=6)
rect(s, Inches(6.6), Inches(1.5), Inches(6.15), Inches(4.9), LIGHT)
txt(s, Inches(6.85), Inches(1.7), Inches(5.6), Inches(0.4), "MODELS COMPARED", size=14, color=NAVY, bold=True)
bullets(s, Inches(6.85), Inches(2.25), Inches(5.6), Inches(2.2), [
    ("Logistic Regression ", "— interpretable, transparent baseline."),
    ("Decision Tree ", "— captures simple non-linear rules."),
    ("Gradient Boosting ", "— ensemble of trees; strongest on tabular telecom data."),
], size=14, gap=10)
txt(s, Inches(6.85), Inches(4.55), Inches(5.6), Inches(1.6),
    "Each model handles class balance and is scored on identical folds, so the comparison "
    "is apples-to-apples.", size=13.5, color=DARK, line_spacing=1.1)
footer_src(s, "Reproducible analysis in accompanying notebook (.ipynb).")

# ============================================================
# SLIDE 9 — Model results
# ============================================================
s = slide(); header(s, 9, "Results: Gradient Boosting Is the Clear Winner", "Model evaluation")
pic(s, ASSET+'04_model_compare.png', Inches(0.5), Inches(1.5), w=Inches(6.0))
pic(s, ASSET+'05_roc.png', Inches(7.0), Inches(1.45), w=Inches(4.4))
rect(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.8), LIGHT)
txt(s, Inches(0.75), Inches(5.97), Inches(12), Inches(0.75),
    "Model: Gradient Boosting Classifier   ·   Metric: ROC-AUC   ·   Score: 0.69 (test)   "
    "—   ~38% better than a random baseline (0.50).",
    size=15, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer_src(s, "5-fold CV: LogReg 0.629, Decision Tree 0.641, Gradient Boosting 0.690. Held-out test AUC: 0.690.")

# ============================================================
# SLIDE 10 — From prediction to action
# ============================================================
s = slide(); header(s, 10, "From Prediction to Action: Who to Call First", "Operationalizing the model")
pic(s, ASSET+'06_lift.png', Inches(0.5), Inches(1.55), w=Inches(8.1))
rect(s, Inches(8.95), Inches(1.7), Inches(3.85), Inches(4.6), LIGHT)
rect(s, Inches(8.95), Inches(1.7), Inches(3.85), Pt(5), RED)
txt(s, Inches(9.2), Inches(2.0), Inches(3.4), Inches(0.5), "READING THE CHART", size=14, color=NAVY, bold=True)
bullets(s, Inches(9.2), Inches(2.6), Inches(3.4), Inches(3.6), [
    ("Top 5%: ", "83% are real churners — 1.67× better than random."),
    ("Top 20%: ", "73% precision, capturing ~29% of all churners."),
    ("Takeaway: ", "a small, focused campaign captures a large share of churn risk cheaply."),
], size=13.5, gap=12)
footer_src(s, "Lift and precision/recall computed on the 20,000-customer held-out test set.")

# ============================================================
# SLIDE 11 — Business proposal
# ============================================================
s = slide(); header(s, 11, "The Proposal: A Monthly Churn-Risk Engine", "Recommended solution")
# 3-step flow
steps = [
    ("1.  SCORE", "Every month, score all customers 0–1 for churn risk using the Gradient Boosting model.", BLUE),
    ("2.  RANK & TARGET", "Send the retention team a ranked list; focus on the top 20% highest-risk customers.", TEAL),
    ("3.  INTERVENE", "Trigger tailored offers — device upgrades, loyalty discounts, proactive support — by risk driver.", GREEN),
]
x = Inches(0.55)
for title_, body, col in steps:
    rect(s, x, Inches(1.6), Inches(3.95), Inches(2.5), LIGHT)
    rect(s, x, Inches(1.6), Inches(3.95), Inches(0.6), col)
    txt(s, x, Inches(1.62), Inches(3.95), Inches(0.55), title_, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(0.2), Inches(2.4), Inches(3.55), Inches(1.6), body, size=14, color=DARK, line_spacing=1.1)
    x = Emu(x + Inches(4.18))
txt(s, Inches(0.55), Inches(4.45), Inches(12), Inches(0.5), "Why it works", size=17, color=NAVY, bold=True)
bullets(s, Inches(0.55), Inches(5.0), Inches(12.2), Inches(2), [
    ("Targets the right people ", "— EDA-validated drivers (old handsets, declining usage) map to concrete offers."),
    ("Fits existing operations ", "— plugs into the current retention team; no new infrastructure needed for the PoC."),
    ("Spends efficiently ", "— contacting 20% instead of 100% cuts campaign cost 5× while keeping most of the upside."),
], size=15, gap=9)

# ============================================================
# SLIDE 12 — Quantified impact
# ============================================================
s = slide(); header(s, 12, "Quantified Impact", "The business case")
pic(s, ASSET+'07_business.png', Inches(0.5), Inches(1.5), w=Inches(7.2))
rect(s, Inches(8.0), Inches(1.55), Inches(4.8), Inches(4.9), LIGHT)
txt(s, Inches(8.25), Inches(1.75), Inches(4.4), Inches(0.4), "TOP-20% CAMPAIGN", size=14, color=NAVY, bold=True)
bullets(s, Inches(8.25), Inches(2.3), Inches(4.4), Inches(3.2), [
    ("20,000 ", "customers contacted"),
    ("73% ", "precision → 14,560 churners reached"),
    ("30% ", "save rate → ~4,368 retained"),
    ("$58.72 ", "avg monthly revenue / customer"),
    ("24-month ", "lifetime-value horizon"),
], size=14, gap=9)
rect(s, Inches(8.0), Inches(5.55), Inches(4.8), Inches(0.9), NAVY)
txt(s, Inches(8.0), Inches(5.6), Inches(4.8), Inches(0.8),
    "$6.0M net benefit  ·  +$2.0M vs. random", size=15.5, color=WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer_src(s, "Assumptions: 30% campaign success, $5/contact, 24-mo LTV, avg revenue from data. Conservative, stated for transparency.")

# ============================================================
# SLIDE 13 — Roadmap / next steps
# ============================================================
s = slide(); header(s, 13, "Next Steps & Risks", "From PoC to production")
txt(s, Inches(0.55), Inches(1.5), Inches(6), Inches(0.5), "Phased roadmap", size=17, color=NAVY, bold=True)
bullets(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.5), [
    ("Phase 1 (0–1 mo): ", "A/B test the top-20% campaign on a hold-out segment to measure true save-rate."),
    ("Phase 2 (1–3 mo): ", "Feed results back; tune offers by churn driver; add SHAP for per-customer explanations."),
    ("Phase 3 (3–6 mo): ", "Automate monthly scoring; try XGBoost/LightGBM & uplift modeling; retrain to counter drift."),
], size=15, gap=12)
txt(s, Inches(6.9), Inches(1.5), Inches(6), Inches(0.5), "Risks & mitigations", size=17, color=NAVY, bold=True)
bullets(s, Inches(6.9), Inches(2.05), Inches(5.9), Inches(4.5), [
    ("Model drift ", "→ scheduled monthly retraining."),
    ("Save-rate uncertainty ", "→ validate via controlled A/B test before scaling."),
    ("Offer cannibalization ", "→ cap discounts; target by risk driver, not blanket."),
    ("Data quality / privacy ", "→ governance review; the model uses only provided, consented data."),
], size=15, gap=12)

# ============================================================
# SLIDE 14 — References
# ============================================================
s = slide(); header(s, 14, "References", "Sources & citations")
refs = [
    "[1] IBISWorld / Statista (2024). U.S. Wireless Telecommunications Carriers — Market Size.",
    "[2] McKinsey & Company. Reducing churn in telecom through advanced analytics.",
    "[3] Harvard Business Review. The Value of Keeping the Right Customers (Gallo, 2014).",
    "[4] Reichheld, F. & Sasser, W. E. Zero Defections: Quality Comes to Services, HBR; Bain & Company.",
    "[5] Dataset: Company A (anonymized telecom) PoC dataset — GCI World 2026 course materials.",
    "[6] Pedregosa et al. (2011). Scikit-learn: Machine Learning in Python. JMLR 12.",
    "[7] Friedman, J. (2001). Greedy Function Approximation: A Gradient Boosting Machine. Annals of Statistics.",
    "[8] Course lecture notebooks, Matsuo–Iwasawa Lab (GCI World 2026 April).",
]
bullets(s, Inches(0.55), Inches(1.6), Inches(12.2), Inches(5.2), refs, size=14, gap=10, bullet='')
txt(s, Inches(0.55), Inches(6.7), Inches(12), Inches(0.4),
    "Generative AI (Claude, Anthropic) was used to assist with structuring and visualization; chat URL listed in the Omnicampus References field.",
    size=11, color=GRAY, italic=True)

# ---- save ----
out_pptx = 'JosselineDanielaChavezCarpio.pptx'
prs.save(out_pptx)
print(f'Saved {out_pptx} with {len(prs.slides._sldIdLst)} slides')
